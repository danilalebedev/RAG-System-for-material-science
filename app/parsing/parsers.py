from __future__ import annotations

import hashlib
import json
import re
import csv
import time
import unicodedata
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from datetime import date, datetime
from io import StringIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

from app.io_utils import safe_filename
from app.settings import paths


MAX_SPREADSHEET_PREVIEW_ROWS_PER_SHEET = 80
MAX_SPREADSHEET_PREVIEW_COLUMNS = 50
MAX_SPREADSHEET_SHEETS = 40
MAX_SPREADSHEET_TEXT_CHARS = 80_000
MAX_SPREADSHEET_TABLE_TEXT_CHARS = 12_000
CSV_EXPORT_PAUSE_EVERY_ROWS = 5000
CSV_EXPORT_PAUSE_SECONDS = 0.02


@dataclass
class ParsedTable:
    table_id: str
    page_or_sheet: str
    rows: list[list[str]]
    text: str


@dataclass
class ParsedDocument:
    local_path: str
    parser: str
    status: str
    title: str = ""
    text: str = ""
    page_count: int = 0
    tables: list[ParsedTable] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)


def file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def parse_document(path: Path) -> ParsedDocument:
    suffix = path.suffix.lower()
    try:
        if looks_like_image(path):
            return parse_image_metadata(path)
        if suffix == ".pdf":
            return parse_pdf(path)
        if suffix == ".docx":
            return parse_docx(path)
        if suffix == ".docm":
            return parse_docm(path)
        if suffix == ".pptx":
            return parse_pptx(path)
        if suffix in {".xlsx", ".xls"}:
            return parse_spreadsheet(path)
        if suffix == ".doc":
            return parse_legacy_doc(path)
        if suffix == ".txt":
            return ParsedDocument(str(path), "text", "ok", text=path.read_text(encoding="utf-8", errors="replace"))
        if suffix in {".gif", ".bmp"}:
            return parse_image_metadata(path)
        return ParsedDocument(str(path), "unsupported", "unsupported", errors=[f"unsupported extension: {suffix}"])
    except Exception as exc:  # noqa: BLE001
        return ParsedDocument(str(path), parser=suffix.lstrip(".") or "unknown", status="failed", errors=[str(exc)])


def parse_pdf(path: Path) -> ParsedDocument:
    doc = fitz.open(path)
    pages: list[str] = []
    metadata = {key: value for key, value in (doc.metadata or {}).items() if value}
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        if text.strip():
            pages.append(f"\n\n--- PAGE {i} ---\n{text.strip()}")
    return ParsedDocument(
        local_path=str(path),
        parser="pymupdf",
        status="ok" if pages else "empty",
        title=metadata.get("title", "") or path.stem,
        text="\n".join(pages),
        page_count=doc.page_count,
        metadata=metadata,
    )


def parse_docx(path: Path) -> ParsedDocument:
    doc = DocxDocument(path)
    parts: list[str] = []
    tables: list[ParsedTable] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for idx, table in enumerate(doc.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        table_text = "\n".join(" | ".join(cell for cell in row if cell) for row in rows)
        if table_text.strip():
            table_id = f"{path.stem}_table_{idx}"
            tables.append(ParsedTable(table_id=table_id, page_or_sheet=f"table_{idx}", rows=rows, text=table_text))
            parts.append(f"\n[TABLE {idx}]\n{table_text}")
    props = doc.core_properties
    metadata = {
        "author": props.author,
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
    }
    text = "\n".join(parts)
    return ParsedDocument(str(path), "python-docx", "ok" if text.strip() else "empty", path.stem, text, 0, tables, [], metadata)


def parse_docm(path: Path) -> ParsedDocument:
    try:
        return parse_docx(path)
    except Exception:
        return parse_word_openxml(path, parser="docm-openxml-fallback")


def parse_word_openxml(path: Path, parser: str) -> ParsedDocument:
    parts: list[str] = []
    tables: list[ParsedTable] = []
    with zipfile.ZipFile(path) as zf:
        document_xml = zf.read("word/document.xml")
    root = ET.fromstring(document_xml)
    namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    for element in root.iter():
        if element.tag == f"{{{namespaces['w']}}}p":
            text = "".join(node.text or "" for node in element.findall(".//w:t", namespaces)).strip()
            if text:
                parts.append(text)
        elif element.tag == f"{{{namespaces['w']}}}tbl":
            rows: list[list[str]] = []
            for row in element.findall(".//w:tr", namespaces):
                cells = [
                    " ".join(node.text or "" for node in cell.findall(".//w:t", namespaces)).strip()
                    for cell in row.findall("./w:tc", namespaces)
                ]
                if any(cells):
                    rows.append(cells)
            if rows:
                table_text = "\n".join(" | ".join(cell for cell in row if cell) for row in rows)
                table_id = f"{path.stem}_table_{len(tables) + 1}"
                tables.append(ParsedTable(table_id=table_id, page_or_sheet=f"table_{len(tables) + 1}", rows=rows, text=table_text))
                parts.append(f"\n[TABLE {len(tables)}]\n{table_text}")

    text = "\n".join(parts)
    return ParsedDocument(
        str(path),
        parser,
        "ok" if text.strip() else "empty",
        path.stem,
        text,
        tables=tables,
        metadata={"extraction_note": "OpenXML fallback for macro-enabled Word document; macros ignored"},
    )


def parse_pptx(path: Path) -> ParsedDocument:
    prs = Presentation(path)
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            note_text = notes_frame.text.strip() if notes_frame is not None and notes_frame.text else ""
            if note_text:
                texts.append(f"NOTES: {note_text}")
        if texts:
            slides.append(f"\n\n--- SLIDE {i} ---\n" + "\n".join(texts))
    text = "\n".join(slides)
    return ParsedDocument(str(path), "python-pptx", "ok" if text.strip() else "empty", path.stem, text, len(prs.slides))


def looks_like_image(path: Path) -> bool:
    try:
        header = path.read_bytes()[:16]
    except OSError:
        return False
    return header.startswith((b"GIF87a", b"GIF89a", b"BM"))


def parse_image_metadata(path: Path) -> ParsedDocument:
    data = path.read_bytes()[:32]
    metadata: dict[str, Any] = {"file_size": path.stat().st_size}
    image_type = "unknown"
    width: int | None = None
    height: int | None = None
    if data.startswith((b"GIF87a", b"GIF89a")) and len(data) >= 10:
        image_type = "gif"
        width = int.from_bytes(data[6:8], "little")
        height = int.from_bytes(data[8:10], "little")
    elif data.startswith(b"BM") and len(data) >= 26:
        image_type = "bmp"
        width = int.from_bytes(data[18:22], "little", signed=True)
        height = abs(int.from_bytes(data[22:26], "little", signed=True))
    metadata.update({"image_type": image_type, "width": width, "height": height, "ocr_performed": False})
    text = (
        f"Image file: {path.name}\n"
        f"Image type: {image_type}\n"
        f"Dimensions: {width or 'unknown'} x {height or 'unknown'}\n"
        "OCR text: not extracted in current parser; use source image for visual table/content inspection."
    )
    return ParsedDocument(str(path), "image-metadata", "ok", path.stem, text, metadata=metadata)


def parse_legacy_doc(path: Path) -> ParsedDocument:
    data = path.read_bytes()
    stripped = data.lstrip()
    if stripped.startswith(b"{\\rtf"):
        text = parse_rtf_text(data)
        return ParsedDocument(str(path), "rtf-fallback", "ok" if text.strip() else "empty", path.stem, text)

    candidates: list[tuple[int, str]] = []
    for offset, encoding in ((0, "utf-16le"), (1, "utf-16le"), (0, "cp1251"), (0, "cp1252")):
        try:
            decoded = data[offset:].decode(encoding, errors="ignore")
        except LookupError:
            continue
        candidates.extend((offset + start, text) for start, text in extract_text_runs(decoded))

    text = dedupe_text_runs(candidates)
    return ParsedDocument(
        str(path),
        "legacy-doc-fallback",
        "ok" if text.strip() else "empty",
        path.stem,
        text,
        metadata={"extraction_note": "best-effort text run extraction; no Microsoft Office automation"},
    )


def parse_rtf_text(data: bytes) -> str:
    raw = data.decode("cp1251", errors="ignore")
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", raw)
    raw = raw.replace("{", " ").replace("}", " ").replace("\\", " ")
    return normalize_text(raw)


def extract_text_runs(text: str, min_chars: int = 24) -> list[tuple[int, str]]:
    runs: list[tuple[int, str]] = []
    start: int | None = None
    buffer: list[str] = []
    for index, char in enumerate(text):
        if is_probable_text_char(char):
            if start is None:
                start = index
            buffer.append(char)
            continue
        if start is not None:
            add_text_run(runs, start, buffer, min_chars)
            start = None
            buffer = []
    if start is not None:
        add_text_run(runs, start, buffer, min_chars)
    return runs


def add_text_run(runs: list[tuple[int, str]], start: int, buffer: list[str], min_chars: int) -> None:
    value = normalize_text("".join(buffer))
    if len(value) < min_chars:
        return
    lowered = value.casefold()
    if any(marker in lowered for marker in ("<?xml", "xmlns:", "theme/theme", "_rels/", "[content_types]")):
        return
    alpha_num_chars = [char.casefold() for char in value if char.isalpha() or char.isdigit()]
    if len(alpha_num_chars) / max(len(value), 1) < 0.35:
        return
    if len(set(alpha_num_chars)) < 8:
        return
    most_common = Counter(alpha_num_chars).most_common(1)
    if most_common and most_common[0][1] / max(len(alpha_num_chars), 1) > 0.35:
        return
    runs.append((start, value))


def is_probable_text_char(char: str) -> bool:
    if char in "\n\r\t ":
        return True
    code = ord(char)
    if 0x20 <= code <= 0x7E:
        return True
    if 0x0400 <= code <= 0x052F:
        return True
    if char in "№°±×÷µ–—«»„“”€₽":
        return True
    category = unicodedata.category(char)
    return category in {"Nd", "Pc", "Pd", "Pe", "Pf", "Pi", "Po", "Ps"}


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def dedupe_text_runs(candidates: list[tuple[int, str]]) -> str:
    lines: list[str] = []
    seen: set[str] = set()
    for _, text in sorted(candidates, key=lambda item: item[0]):
        normalized = normalize_text(text)
        key = normalized.casefold()
        if key in seen:
            continue
        seen.add(key)
        lines.append(normalized)
    return "\n".join(lines)


def parse_spreadsheet(path: Path) -> ParsedDocument:
    if path.suffix.lower() == ".xls":
        return parse_xls_to_csv(path)
    return parse_xlsx_to_csv(path)


def parse_xlsx_to_csv(path: Path) -> ParsedDocument:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    tables: list[ParsedTable] = []
    metadata = spreadsheet_metadata(path, len(workbook.worksheets), parser="openpyxl-csv-export")
    try:
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            if sheet_index > MAX_SPREADSHEET_SHEETS:
                metadata["preview_truncated"] = True
                parts.append(f"\n[SPREADSHEET PREVIEW TRUNCATED]\nSkipped preview after {MAX_SPREADSHEET_SHEETS} sheets.")
                break
            csv_path, preview_rows, row_count, column_count = export_xlsx_sheet(path, worksheet, sheet_index)
            append_spreadsheet_sheet(
                path=path,
                sheet_index=sheet_index,
                sheet_name=worksheet.title,
                csv_path=csv_path,
                row_count=row_count,
                column_count=column_count,
                preview_rows=preview_rows,
                parts=parts,
                tables=tables,
                metadata=metadata,
            )
    finally:
        workbook.close()
    return spreadsheet_document(path, "openpyxl-csv-export", parts, tables, metadata)


def parse_xls_to_csv(path: Path) -> ParsedDocument:
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    sheet_names = workbook.sheet_names()
    parts: list[str] = []
    tables: list[ParsedTable] = []
    metadata = spreadsheet_metadata(path, len(sheet_names), parser="xlrd-csv-export")
    try:
        for sheet_index, sheet_name in enumerate(sheet_names, start=1):
            if sheet_index > MAX_SPREADSHEET_SHEETS:
                metadata["preview_truncated"] = True
                parts.append(f"\n[SPREADSHEET PREVIEW TRUNCATED]\nSkipped preview after {MAX_SPREADSHEET_SHEETS} sheets.")
                break
            sheet = workbook.sheet_by_name(sheet_name)
            csv_path, preview_rows, row_count, column_count = export_xls_sheet(path, workbook, sheet, sheet_index)
            append_spreadsheet_sheet(
                path=path,
                sheet_index=sheet_index,
                sheet_name=sheet_name,
                csv_path=csv_path,
                row_count=row_count,
                column_count=column_count,
                preview_rows=preview_rows,
                parts=parts,
                tables=tables,
                metadata=metadata,
            )
            workbook.unload_sheet(sheet_name)
    finally:
        workbook.release_resources()
    return spreadsheet_document(path, "xlrd-csv-export", parts, tables, metadata)


def spreadsheet_metadata(path: Path, sheet_count: int, parser: str) -> dict[str, Any]:
    return {
        "sheet_count": sheet_count,
        "sheets": [],
        "csv_export_dir": str(spreadsheet_export_dir(path)),
        "csv_export_complete": True,
        "parser": parser,
        "preview_row_limit_per_sheet": MAX_SPREADSHEET_PREVIEW_ROWS_PER_SHEET,
        "preview_column_limit_per_sheet": MAX_SPREADSHEET_PREVIEW_COLUMNS,
        "preview_sheet_limit": MAX_SPREADSHEET_SHEETS,
        "preview_text_char_limit": MAX_SPREADSHEET_TEXT_CHARS,
        "preview_truncated": sheet_count > MAX_SPREADSHEET_SHEETS,
    }


def spreadsheet_export_dir(path: Path) -> Path:
    return paths().spreadsheet_csv_dir / f"{stable_id(str(path))}__{safe_filename(path.stem)}"


def sheet_csv_path(path: Path, sheet_index: int, sheet_name: str) -> Path:
    return spreadsheet_export_dir(path) / f"{sheet_index:03d}__{safe_filename(sheet_name)}.csv"


def export_xlsx_sheet(path: Path, worksheet: Any, sheet_index: int) -> tuple[Path, list[list[str]], int, int]:
    csv_path = sheet_csv_path(path, sheet_index, worksheet.title)
    csv_path.parent.mkdir(parents=True, exist_ok=True)
    preview_rows: list[list[str]] = []
    row_count = 0
    column_count = int(worksheet.max_column or 0)
    with csv_path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        for row in worksheet.iter_rows(values_only=True):
            values = [format_cell(value) for value in row]
            writer.writerow(values)
            row_count += 1
            column_count = max(column_count, len(values))
            if len(preview_rows) < MAX_SPREADSHEET_PREVIEW_ROWS_PER_SHEET:
                preview_rows.append(values[:MAX_SPREADSHEET_PREVIEW_COLUMNS])
            if row_count % CSV_EXPORT_PAUSE_EVERY_ROWS == 0:
                time.sleep(CSV_EXPORT_PAUSE_SECONDS)
    return csv_path, preview_rows, row_count, column_count


def export_xls_sheet(path: Path, workbook: Any, sheet: Any, sheet_index: int) -> tuple[Path, list[list[str]], int, int]:
    csv_path = sheet_csv_path(path, sheet_index, sheet.name)
    csv_path.parent.mkdir(parents=True, exist_ok=True)
    preview_rows: list[list[str]] = []
    with csv_path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        for row_index in range(sheet.nrows):
            values = [format_xls_cell(workbook, sheet.cell(row_index, column_index)) for column_index in range(sheet.ncols)]
            writer.writerow(values)
            if len(preview_rows) < MAX_SPREADSHEET_PREVIEW_ROWS_PER_SHEET:
                preview_rows.append(values[:MAX_SPREADSHEET_PREVIEW_COLUMNS])
            if (row_index + 1) % CSV_EXPORT_PAUSE_EVERY_ROWS == 0:
                time.sleep(CSV_EXPORT_PAUSE_SECONDS)
    return csv_path, preview_rows, sheet.nrows, sheet.ncols


def append_spreadsheet_sheet(
    *,
    path: Path,
    sheet_index: int,
    sheet_name: str,
    csv_path: Path,
    row_count: int,
    column_count: int,
    preview_rows: list[list[str]],
    parts: list[str],
    tables: list[ParsedTable],
    metadata: dict[str, Any],
) -> None:
    preview_text = rows_to_csv(preview_rows)
    if len(preview_text) > MAX_SPREADSHEET_TABLE_TEXT_CHARS:
        metadata["preview_truncated"] = True
        preview_text = preview_text[:MAX_SPREADSHEET_TABLE_TEXT_CHARS].rstrip() + "\n...[preview truncated]"

    sheet_meta = {
        "sheet_index": sheet_index,
        "sheet_name": sheet_name,
        "rows": row_count,
        "columns": column_count,
        "csv_path": str(csv_path),
        "csv_size": csv_path.stat().st_size if csv_path.exists() else 0,
        "preview_rows": len(preview_rows),
        "preview_columns": min(column_count, MAX_SPREADSHEET_PREVIEW_COLUMNS),
    }
    metadata["sheets"].append(sheet_meta)

    table_id = f"{path.stem}_{sheet_name}".replace(" ", "_")
    tables.append(ParsedTable(table_id=table_id, page_or_sheet=sheet_name, rows=preview_rows, text=preview_text))

    sheet_part = (
        f"\n[SHEET {sheet_index}: {sheet_name}]\n"
        f"rows={row_count}, columns={column_count}, csv_path={csv_path}\n"
        f"[PREVIEW]\n{preview_text}"
    )
    existing_chars = sum(len(part) for part in parts)
    if existing_chars + len(sheet_part) > MAX_SPREADSHEET_TEXT_CHARS:
        metadata["preview_truncated"] = True
        remaining = max(0, MAX_SPREADSHEET_TEXT_CHARS - existing_chars)
        if remaining:
            parts.append(sheet_part[:remaining].rstrip())
        parts.append("\n[SPREADSHEET PREVIEW TRUNCATED]\nPreview text char limit reached; full data is in CSV exports.")
        return
    parts.append(sheet_part)


def rows_to_csv(rows: list[list[str]]) -> str:
    buffer = StringIO()
    writer = csv.writer(buffer, lineterminator="\n")
    writer.writerows(rows)
    return buffer.getvalue()


def format_xls_cell(workbook: Any, cell: Any) -> str:
    import xlrd

    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            return xlrd.xldate.xldate_as_datetime(cell.value, workbook.datemode).isoformat(sep=" ")
        except Exception:  # noqa: BLE001
            return format_cell(cell.value)
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if bool(cell.value) else "FALSE"
    return format_cell(cell.value)


def format_cell(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.isoformat(sep=" ")
    if isinstance(value, date):
        return value.isoformat()
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def spreadsheet_document(
    path: Path,
    parser: str,
    parts: list[str],
    tables: list[ParsedTable],
    metadata: dict[str, Any],
) -> ParsedDocument:
    return ParsedDocument(
        str(path),
        parser,
        "ok" if metadata.get("sheets") else "empty",
        path.stem,
        "\n".join(parts),
        int(metadata.get("sheet_count", 0)),
        tables,
        metadata=metadata,
    )


def parsed_document_to_row(doc: ParsedDocument, extra: dict[str, Any] | None = None) -> dict[str, Any]:
    path = Path(doc.local_path)
    row = {
        "doc_id": stable_id(str(path)),
        "local_path": doc.local_path,
        "file_name": path.name,
        "extension": path.suffix.lower(),
        "parser": doc.parser,
        "status": doc.status,
        "title": doc.title,
        "page_count": doc.page_count,
        "text_chars": len(doc.text),
        "table_count": len(doc.tables),
        "errors": "; ".join(doc.errors),
        "metadata_json": json.dumps(doc.metadata, ensure_ascii=False, default=str),
    }
    if extra:
        row.update(extra)
    return row


def stable_id(value: str) -> str:
    return hashlib.sha1(value.encode("utf-8", errors="ignore")).hexdigest()[:16]
